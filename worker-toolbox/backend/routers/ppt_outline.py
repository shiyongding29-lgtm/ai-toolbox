"""PPT / HTML 大纲生成 + PPTX / HTML 下载 API。"""
from io import BytesIO
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from backend.services.llm_service import llm_service
from backend.services.prompt_library import PPT_OUTLINE_SYSTEM

router = APIRouter(prefix="/api/ppt-outline", tags=["ppt-outline"])


class LlmReq(BaseModel):
    text: str = ""
    extra_context: str = ""


@router.post("/run")
async def ppt_outline_run(req: LlmReq):
    prompt = f"Topic: {req.text}"
    if req.extra_context:
        prompt += f"\n\nRequirements: {req.extra_context}"
    result = llm_service.complete(PPT_OUTLINE_SYSTEM, prompt)
    return {"code": 0, "data": {"result": result}}


class GeneratePptxRequest(BaseModel):
    outline: str = ""
    slide_count: int = 12


def generate_pptx_bytes(outline: str, slide_count: int = 12) -> BytesIO:
    """根据大纲生成 PPTX 文件，返回 BytesIO buffer。自动从大纲中提取风格配色。"""
    import re

    # 从大纲中提取配色方案
    palette_match = re.search(r'Palette:\s*([#0-9a-fA-Fx,\s]+)', outline)
    colors = []
    if palette_match:
        hex_strs = re.findall(r'#[0-9a-fA-F]{6}', palette_match.group(1))
        colors = [rgb_from_hex(h) for h in hex_strs]

    # 默认配色：专业蓝
    if not colors:
        colors = [RGBColor(0x1d, 0x42, 0x8a), RGBColor(0xff, 0xff, 0xff), RGBColor(0xd4, 0xa8, 0x53)]
    while len(colors) < 3:
        colors.append(RGBColor(0x66, 0x66, 0x66))

    primary = colors[0]   # 标题色
    bg_light = colors[1] if len(colors) > 1 else RGBColor(0xff, 0xff, 0xff)
    accent = colors[2] if len(colors) > 2 else RGBColor(0x25, 0x63, 0xeb)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    lines = outline.split('\n')
    slides = []
    current_slide = None
    current_bullets = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        # Skip style header lines
        if stripped.startswith('## Style:') or stripped.startswith('## Design System') or stripped.startswith('## Typography') or stripped.startswith('## Color Application') or stripped.startswith('## Chart/Diagram') or stripped.startswith('## Image Style') or stripped.startswith('## Animation'):
            continue
        if stripped.startswith('- **Layout**') or stripped.startswith('- **Visual**'):
            continue
        if stripped.startswith('## ') or stripped.startswith('# '):
            if current_slide:
                current_slide['bullets'] = current_bullets
                slides.append(current_slide)
            current_slide = {'title': stripped.lstrip('#').strip(), 'bullets': []}
            current_bullets = []
        elif stripped.startswith('- ') or stripped.startswith('* '):
            current_bullets.append(stripped.lstrip('-* ').strip())
        elif current_slide is not None:
            if not current_slide['title']:
                current_slide['title'] = stripped[:80]
            else:
                current_bullets.append(stripped)

    if current_slide:
        current_slide['bullets'] = current_bullets
        slides.append(current_slide)

    if not slides:
        for line in lines:
            if line.strip():
                slides.append({'title': line.strip()[:100], 'bullets': []})

    slides = slides[:slide_count]

    for i, slide_data in enumerate(slides):
        # Use Title and Content layout for all slides (most compatible)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Add a colored accent bar at the top
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = primary
        title.text_frame.paragraphs[0].font.bold = True

        # Bullets
        if slide_data['bullets']:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for j, bullet in enumerate(slide_data['bullets']):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_after = Pt(10)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def rgb_from_hex(hex_str: str) -> RGBColor:
    """将 #rrggbb 转为 python-pptx RGBColor。"""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@router.post("/generate-pptx")
async def generate_pptx(req: GeneratePptxRequest):
    """根据大纲生成 PPTX 文件并下载。"""
    buffer = generate_pptx_bytes(req.outline, req.slide_count)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )


class GenerateHtmlRequest(BaseModel):
    outline: str = ""


def generate_html_string(outline: str) -> str:
    """根据大纲生成 reveal.js HTML 字符串，带 fallback 确保始终可打开。"""
    import re

    # 尝试让 LLM 生成 HTML
    try:
        resp = llm_service.complete(
            "Output ONLY valid HTML code, nothing else. No markdown fences, no explanations.",
            f"Create a complete reveal.js HTML presentation. Use CDN reveal.js@5.0.0. Choose theme based on the outline style (night for dark/tech, simple for minimal, league for business). Each ## heading is a <section>. Use the color palette from the outline. Include Reveal.initialize(). Start with <!DOCTYPE html>:\n{outline}",
        )
        # 提取 HTML
        html_match = re.search(r'(<!DOCTYPE[\s\S]*)', resp)
        if html_match:
            html = html_match.group(1)
            # 验证基本结构
            if '</html>' in html and '<section>' in html:
                return html
    except Exception:
        pass

    # Fallback: 从大纲手动构建 HTML
    return _build_html_fallback(outline)


def _build_html_fallback(outline: str) -> str:
    """从大纲构建基础 reveal.js HTML，确保始终能打开。"""
    import re

    # 尝试提取配色
    palette_match = re.search(r'Palette:\s*([#0-9a-fA-Fx,\s]+)', outline)
    bg_color = "#0a0e27"
    text_color = "#e0e0e0"
    accent_color = "#00e5ff"
    if palette_match:
        hex_strs = re.findall(r'#[0-9a-fA-F]{6}', palette_match.group(1))
        if len(hex_strs) >= 1:
            bg_color = hex_strs[0]
        if len(hex_strs) >= 2:
            text_color = hex_strs[1]
        if len(hex_strs) >= 3:
            accent_color = hex_strs[2]

    # 判断是否深色主题
    is_dark = _is_dark_color(bg_color)
    reveal_theme = "night" if is_dark else "simple"
    slide_bg = bg_color
    font_color = "#ffffff" if is_dark else "#1a1a1a"

    # 解析大纲为 slides
    lines = outline.split('\n')
    sections = []
    current_title = ""
    current_bullets = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith('## Style:') or stripped.startswith('## Design Specs'):
            continue
        if stripped.startswith('## '):
            if current_title:
                sections.append((current_title, current_bullets))
            current_title = stripped.lstrip('#').strip()
            current_bullets = []
        elif stripped.startswith('# '):
            current_title = stripped.lstrip('#').strip()
            current_bullets = []
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullet = stripped.lstrip('-* ').strip()
            if not bullet.startswith('[') and not bullet.startswith('**Layout**') and not bullet.startswith('**Visual**'):
                current_bullets.append(bullet)

    if current_title:
        sections.append((current_title, current_bullets))

    if not sections:
        sections = [("Presentation", ["No content generated"])]

    # 构建 HTML
    slides_html = ""
    for title, bullets in sections:
        slides_html += f'<section>\n  <h2>{_escape_html(title)}</h2>\n'
        if bullets:
            slides_html += '  <ul>\n'
            for b in bullets:
                slides_html += f'    <li>{_escape_html(b)}</li>\n'
            slides_html += '  </ul>\n'
        slides_html += '</section>\n'

    return f"""<!DOCTYPE html>
<html lang="zh">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Presentation</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.0.0/dist/reset.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.0.0/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.0.0/dist/theme/{reveal_theme}.css">
<style>
  :root {{
    --r-background-color: {slide_bg};
    --r-main-color: {font_color};
    --r-heading-color: {accent_color};
    --r-link-color: {accent_color};
  }}
  .reveal h2 {{ color: {accent_color}; text-shadow: 0 0 20px {accent_color}44; }}
  .reveal li {{ margin: 0.5em 0; font-size: 0.85em; }}
  .reveal section {{ padding: 40px; }}
</style>
</head>
<body>
<div class="reveal">
  <div class="slides">
{slides_html}
  </div>
</div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@5.0.0/dist/reveal.js"></script>
<script>
  Reveal.initialize({{ hash: true, slideNumber: true, transition: 'slide', transitionSpeed: 'default' }});
</script>
</body>
</html>"""


def _escape_html(text: str) -> str:
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')


def _is_dark_color(hex_color: str) -> bool:
    """判断 hex 颜色是否为深色。"""
    h = hex_color.lstrip('#')
    if len(h) < 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128


@router.post("/generate-html")
async def generate_html(req: GenerateHtmlRequest):
    """根据大纲生成 reveal.js HTML 演示文稿。"""
    html = generate_html_string(req.outline)
    buffer = BytesIO(html.encode('utf-8'))
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="text/html",
        headers={"Content-Disposition": "attachment; filename=presentation.html"},
    )
